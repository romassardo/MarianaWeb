"""
Script para extraer imágenes de una presentación PowerPoint
y optimizarlas para su uso en un sitio web
"""

import os
import sys
from pptx import Presentation
from PIL import Image
import io
import re

def sanitize_filename(text):
    """Convierte un texto a un nombre de archivo válido"""
    text = re.sub(r'[^\w\s-]', '', text).strip().lower()
    text = re.sub(r'[-\s]+', '-', text)
    return text

def extract_images_from_pptx(pptx_path, output_folder, optimize=True):
    """Extrae las imágenes de un archivo PowerPoint y las guarda en la carpeta especificada"""
    # Crear carpeta de salida si no existe
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # Abrir la presentación
    presentation = Presentation(pptx_path)
    
    # Contador para las imágenes sin nombre
    img_count = 1
    extracted_images = []
    
    # Iterar por cada diapositiva
    for i, slide in enumerate(presentation.slides):
        slide_title = f"slide_{i+1}"
        
        # Intentar obtener el título de la diapositiva si existe
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    slide_title = sanitize_filename(shape.text)
                    break
        
        # Extraer las imágenes de la diapositiva
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                # Obtener los datos de la imagen
                image_stream = io.BytesIO(shape.image.blob)
                img = Image.open(image_stream)
                
                # Nombrar la imagen utilizando el título de la diapositiva
                img_filename = f"obra{img_count}_{slide_title}.jpg"
                img_path = os.path.join(output_folder, img_filename)
                
                # Optimizar la imagen para web
                if optimize:
                    # Convertir a RGB si es necesario
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # Redimensionar si es muy grande
                    max_size = 1600
                    if max(img.width, img.height) > max_size:
                        if img.width > img.height:
                            new_width = max_size
                            new_height = int(img.height * (max_size / img.width))
                        else:
                            new_height = max_size
                            new_width = int(img.width * (max_size / img.height))
                        img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    
                    # Guardar con calidad optimizada
                    img.save(img_path, "JPEG", quality=85, optimize=True)
                else:
                    # Guardar sin optimizar
                    img.save(img_path)
                
                # Crear también una versión en miniatura
                thumb_folder = os.path.join(output_folder, "thumbnails")
                if not os.path.exists(thumb_folder):
                    os.makedirs(thumb_folder)
                
                thumb_size = 400
                if max(img.width, img.height) > thumb_size:
                    if img.width > img.height:
                        new_width = thumb_size
                        new_height = int(img.height * (thumb_size / img.width))
                    else:
                        new_height = thumb_size
                        new_width = int(img.width * (thumb_size / img.height))
                    thumb = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                else:
                    thumb = img
                
                thumb_path = os.path.join(thumb_folder, img_filename)
                thumb.save(thumb_path, "JPEG", quality=85, optimize=True)
                
                # Guardar información sobre la imagen extraída
                extracted_images.append({
                    "original_slide": i+1,
                    "filename": img_filename,
                    "path": img_path,
                    "thumbnail": thumb_path
                })
                
                print(f"Imagen extraída: {img_filename}")
                img_count += 1
    
    return extracted_images

def create_simple_html_gallery(images, output_path):
    """Crea una galería HTML simple para visualizar las imágenes extraídas"""
    html = """<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Imágenes extraídas del PowerPoint</title>
    <style>
        body { font-family: Arial, sans-serif; max-width: 1200px; margin: 0 auto; padding: 20px; }
        h1 { text-align: center; }
        .gallery { display: grid; grid-template-columns: repeat(auto-fill, minmax(250px, 1fr)); gap: 20px; }
        .image-item { border: 1px solid #ddd; border-radius: 4px; overflow: hidden; }
        .image-item img { width: 100%; height: auto; display: block; }
        .image-info { padding: 10px; }
        .image-info p { margin: 5px 0; }
    </style>
</head>
<body>
    <h1>Imágenes extraídas del PowerPoint</h1>
    <div class="gallery">
"""
    
    for img in images:
        html += f"""
        <div class="image-item">
            <a href="{img['path']}" target="_blank">
                <img src="{img['thumbnail']}" alt="{img['filename']}">
            </a>
            <div class="image-info">
                <p><strong>Nombre:</strong> {img['filename']}</p>
                <p><strong>Diapositiva:</strong> {img['original_slide']}</p>
            </div>
        </div>
"""
    
    html += """
    </div>
</body>
</html>
"""
    
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html)
    
    print(f"Galería HTML creada en: {output_path}")

if __name__ == "__main__":
    # Rutas del proyecto
    base_path = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(base_path, "Docs", "PRESENTACION MT 2020.pptx")
    output_folder = os.path.join(base_path, "img", "obras")
    
    # Extraer imágenes
    print(f"Extrayendo imágenes de: {pptx_path}")
    print(f"Guardando en: {output_folder}")
    
    try:
        extracted_images = extract_images_from_pptx(pptx_path, output_folder)
        
        # Crear galería HTML para visualizar las imágenes extraídas
        gallery_path = os.path.join(base_path, "extracted_images_gallery.html")
        create_simple_html_gallery(extracted_images, gallery_path)
        
        print(f"Se han extraído {len(extracted_images)} imágenes correctamente.")
        print(f"Puedes ver las imágenes en el navegador abriendo: {gallery_path}")
    except Exception as e:
        print(f"Error al extraer las imágenes: {e}")
        sys.exit(1)
